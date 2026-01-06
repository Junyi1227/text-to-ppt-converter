#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文字轉 PowerPoint 轉換器 - GUI 版本
Text to PowerPoint Converter - GUI Version

適合打包成 Windows .exe 執行檔
"""

import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys


class TextToPPTConverter:
    """文字轉 PPT 轉換器核心"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
        # 樣式設定
        self.title_bg_color = RGBColor(230, 240, 255)
        self.content_bg_color = RGBColor(245, 245, 245)
        self.font_name = "微軟正黑體"
    
    def create_title_slide(self, title_text):
        """建立主題投影片（## 標記）"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 設定背景顏色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.title_bg_color
        
        # 新增標題文字框
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title_text
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.name = self.font_name
        
        return slide
    
    def create_content_slide(self, title_text):
        """建立內文投影片（# 標記）"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 設定背景顏色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.content_bg_color
        
        # 新增標題文字框
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        text_frame = title_box.text_frame
        text_frame.text = title_text
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.name = self.font_name
        
        # 新增內容文字框
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5)
        )
        
        slide.content_textbox = content_box
        return slide
    
    def add_content_to_slide(self, slide, content_text):
        """新增內容到投影片"""
        if not hasattr(slide, 'content_textbox'):
            return
        
        text_frame = slide.content_textbox.text_frame
        
        if text_frame.text == "":
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = content_text
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = self.font_name
    
    def convert_from_text(self, text_content):
        """從文字內容轉換為 PPT"""
        lines = text_content.split('\n')
        current_slide = None
        
        for line in lines:
            line = line.strip()
            
            if line.startswith('##'):
                title = line[2:].strip()
                current_slide = self.create_title_slide(title)
            elif line.startswith('#'):
                title = line[1:].strip()
                current_slide = self.create_content_slide(title)
            elif line and current_slide:
                self.add_content_to_slide(current_slide, line)
        
        return len(self.prs.slides)
    
    def save(self, output_file):
        """儲存 PPT 檔案"""
        self.prs.save(output_file)


class TextToPPTApp:
    """GUI 應用程式"""
    
    def __init__(self, root):
        self.root = root
        self.root.title("文字轉 PowerPoint 工具")
        self.root.geometry("700x600")
        
        # 設定圖示（如果有的話）
        try:
            # 打包時可以包含 icon.ico
            self.root.iconbitmap("icon.ico")
        except:
            pass
        
        self.setup_ui()
    
    def setup_ui(self):
        """設定使用者介面"""
        
        # 標題
        title_label = tk.Label(
            self.root,
            text="文字轉 PowerPoint 工具",
            font=("微軟正黑體", 16, "bold")
        )
        title_label.pack(pady=10)
        
        # 說明文字
        info_text = (
            "文字格式說明：\n"
            "• ## 開頭：建立主題頁面（藍色背景）\n"
            "• # 開頭：建立內文頁面（灰色背景）\n"
            "• 一般文字：加入到前一張投影片"
        )
        info_label = tk.Label(
            self.root,
            text=info_text,
            font=("微軟正黑體", 10),
            justify=tk.LEFT,
            bg="#f0f0f0",
            padx=10,
            pady=10
        )
        info_label.pack(pady=5, padx=20, fill=tk.X)
        
        # 文字輸入區
        input_frame = tk.LabelFrame(
            self.root,
            text="輸入文字內容",
            font=("微軟正黑體", 11, "bold"),
            padx=10,
            pady=10
        )
        input_frame.pack(pady=10, padx=20, fill=tk.BOTH, expand=True)
        
        self.text_area = scrolledtext.ScrolledText(
            input_frame,
            font=("微軟正黑體", 10),
            wrap=tk.WORD,
            width=60,
            height=15
        )
        self.text_area.pack(fill=tk.BOTH, expand=True)
        
        # 預設範例文字
        default_text = """##歡迎使用文字轉 PPT 工具
這是第一張主題投影片

#什麼是這個工具？
自動將文字轉換成 PowerPoint 簡報
支援兩種投影片格式
使用 ## 建立主題頁面
使用 # 建立內文頁面

##開始使用
修改左側文字，然後點擊「轉換為 PPT」按鈕！"""
        
        self.text_area.insert("1.0", default_text)
        
        # 按鈕區
        button_frame = tk.Frame(self.root)
        button_frame.pack(pady=10)
        
        # 從檔案載入按鈕
        load_button = tk.Button(
            button_frame,
            text="📂 載入文字檔",
            font=("微軟正黑體", 11),
            command=self.load_file,
            width=15,
            bg="#e3f2fd"
        )
        load_button.pack(side=tk.LEFT, padx=5)
        
        # 轉換按鈕
        convert_button = tk.Button(
            button_frame,
            text="🎨 轉換為 PPT",
            font=("微軟正黑體", 11, "bold"),
            command=self.convert_to_ppt,
            width=15,
            bg="#c8e6c9"
        )
        convert_button.pack(side=tk.LEFT, padx=5)
        
        # 清除按鈕
        clear_button = tk.Button(
            button_frame,
            text="🗑️ 清除",
            font=("微軟正黑體", 11),
            command=self.clear_text,
            width=15,
            bg="#ffccbc"
        )
        clear_button.pack(side=tk.LEFT, padx=5)
        
        # 狀態列
        self.status_label = tk.Label(
            self.root,
            text="就緒",
            font=("微軟正黑體", 9),
            relief=tk.SUNKEN,
            anchor=tk.W
        )
        self.status_label.pack(side=tk.BOTTOM, fill=tk.X)
    
    def load_file(self):
        """載入文字檔案"""
        file_path = filedialog.askopenfilename(
            title="選擇文字檔案",
            filetypes=[("文字檔案", "*.txt"), ("所有檔案", "*.*")]
        )
        
        if file_path:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                self.text_area.delete("1.0", tk.END)
                self.text_area.insert("1.0", content)
                self.status_label.config(text=f"已載入：{os.path.basename(file_path)}")
            except Exception as e:
                messagebox.showerror("錯誤", f"無法讀取檔案：{str(e)}")
    
    def clear_text(self):
        """清除文字"""
        if messagebox.askyesno("確認", "確定要清除所有文字嗎？"):
            self.text_area.delete("1.0", tk.END)
            self.status_label.config(text="已清除")
    
    def convert_to_ppt(self):
        """轉換為 PowerPoint"""
        text_content = self.text_area.get("1.0", tk.END).strip()
        
        if not text_content:
            messagebox.showwarning("警告", "請輸入文字內容！")
            return
        
        # 選擇儲存位置
        output_file = filedialog.asksaveasfilename(
            title="儲存 PowerPoint 檔案",
            defaultextension=".pptx",
            filetypes=[("PowerPoint 檔案", "*.pptx"), ("所有檔案", "*.*")]
        )
        
        if not output_file:
            return
        
        try:
            self.status_label.config(text="正在轉換...")
            self.root.update()
            
            # 執行轉換
            converter = TextToPPTConverter()
            slide_count = converter.convert_from_text(text_content)
            converter.save(output_file)
            
            self.status_label.config(text=f"完成！已建立 {slide_count} 張投影片")
            
            # 顯示成功訊息
            result = messagebox.showinfo(
                "成功",
                f"✅ 成功建立 PowerPoint 檔案！\n\n"
                f"📊 總共建立 {slide_count} 張投影片\n"
                f"📁 儲存位置：{output_file}\n\n"
                f"是否開啟檔案所在資料夾？"
            )
            
            # 詢問是否開啟資料夾
            if messagebox.askyesno("開啟資料夾", "要開啟檔案所在的資料夾嗎？"):
                import subprocess
                folder_path = os.path.dirname(output_file)
                if sys.platform == 'win32':
                    os.startfile(folder_path)
                elif sys.platform == 'darwin':
                    subprocess.Popen(['open', folder_path])
                else:
                    subprocess.Popen(['xdg-open', folder_path])
        
        except Exception as e:
            self.status_label.config(text="轉換失敗")
            messagebox.showerror("錯誤", f"轉換時發生錯誤：{str(e)}")


def main():
    """主程式"""
    root = tk.Tk()
    app = TextToPPTApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
