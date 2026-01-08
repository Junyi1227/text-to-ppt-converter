#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
使用模板生成 PPT
根據配置檔和藍色文字清單，生成完整的 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import sys
import os
import re
from copy import deepcopy


class PPTGenerator:
    """PPT 生成器"""
    
    def __init__(self, template_path):
        """初始化，載入模板"""
        self.template_path = template_path
        self.template = Presentation(template_path)
        self.new_prs = Presentation(template_path)
        
        # 清空新簡報（我們會手動複製需要的頁面）
        while len(self.new_prs.slides) > 0:
            rId = self.new_prs.slides._sldIdLst[0].rId
            self.new_prs.part.drop_rel(rId)
            del self.new_prs.slides._sldIdLst[0]
        
        self.config = {}
        self.blue_texts = []
    
    def load_config(self, config_path):
        """讀取配置檔（支援兩種格式：KEY: VALUE 或 KEY=VALUE）"""
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.strip()
                    # 跳過註解和空行
                    if not line or line.startswith('#'):
                        continue
                    
                    # 解析 KEY: VALUE 格式（舊格式）
                    if ':' in line:
                        key, value = line.split(':', 1)
                        self.config[key.strip()] = value.strip()
                    # 解析 KEY=VALUE 格式（新格式）
                    elif '=' in line:
                        key, value = line.split('=', 1)
                        self.config[key.strip()] = value.strip()
            
            print(f"✅ 讀取配置檔：{len(self.config)} 個設定")
            return True
        
        except Exception as e:
            print(f"❌ 讀取配置檔錯誤：{e}")
            return False
    
    def load_blue_texts(self, blue_text_path):
        """讀取藍色文字清單（支援新格式：含 [變數] 區塊）"""
        try:
            with open(blue_text_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 檢查是否為新格式（含 [變數] 區塊）
            if '[變數]' in content and '[變數結束]' in content:
                # 分離變數區和內容區
                parts = content.split('[變數結束]')
                if len(parts) >= 2:
                    var_section = parts[0].replace('[變數]', '').strip()
                    content_section = parts[1].strip()
                    
                    # 解析變數區（支援 = 分隔）
                    for line in var_section.split('\n'):
                        line = line.strip()
                        if '=' in line:
                            key, value = line.split('=', 1)
                            key = key.strip()
                            value = value.strip()
                            
                            # 映射到舊的 config key 名稱
                            key_mapping = {
                                '日期': 'DATE',
                                '禮拜類型': 'SERVICE_TYPE',
                                '主題': 'TITLE',
                                '經文章節': 'VERSE_REFS',
                                '經文1': 'VERSE_1',
                                '經文2': 'VERSE_2',
                            }
                            
                            config_key = key_mapping.get(key, key)
                            self.config[config_key] = value
                            
                            # 解析經文格式（提取章節和內容）
                            if key == '經文1' or key == '經文2':
                                verse_num = '1' if key == '經文1' else '2'
                                # 格式：〈章節〉內容。
                                import re
                                match = re.match(r'^[〈<]([^〉>]+)[〉>]\s*(.+)$', value)
                                if match:
                                    self.config[f'VERSE_REF_{verse_num}'] = match.group(1).strip()
                                    self.config[f'VERSE_TEXT_{verse_num}'] = match.group(2).strip()
                    
                    # 用空行分隔內容段落
                    self.blue_texts = [para.strip() for para in content_section.split('\n\n') if para.strip()]
                    
                    print(f"✅ 讀取新格式：{len(self.config)} 個變數，{len(self.blue_texts)} 段內容")
                    return True
            
            # 舊格式：直接用空行分隔
            self.blue_texts = [para.strip() for para in content.split('\n\n') if para.strip()]
            print(f"✅ 讀取藍色文字：{len(self.blue_texts)} 段")
            return True
        
        except Exception as e:
            print(f"❌ 讀取藍色文字錯誤：{e}")
            return False
    
    def copy_slide(self, slide_index):
        """複製指定投影片到新簡報"""
        source_slide = self.template.slides[slide_index]
        
        # 複製投影片版面配置
        slide_layout = source_slide.slide_layout
        new_slide = self.new_prs.slides.add_slide(slide_layout)
        
        # 背景保持使用版面配置的背景，不手動複製
        # （因為模板已經包含背景設定）
        
        # 複製所有形狀
        for shape in source_slide.shapes:
            self._copy_shape(shape, new_slide)
        
        return new_slide
    
    def _copy_shape(self, source_shape, target_slide):
        """複製形狀到目標投影片"""
        try:
            # 複製文字框
            if hasattr(source_shape, "text_frame"):
                new_shape = target_slide.shapes.add_textbox(
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height
                )
                
                # 複製文字內容和格式
                for paragraph in source_shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        p = new_shape.text_frame.add_paragraph() if new_shape.text_frame.text else new_shape.text_frame.paragraphs[0]
                        p.text = paragraph.text
                        p.alignment = paragraph.alignment
                        
                        # 複製字體格式
                        if paragraph.runs:
                            source_run = paragraph.runs[0]
                            for run in p.runs:
                                if source_run.font.size:
                                    run.font.size = source_run.font.size
                                if source_run.font.bold:
                                    run.font.bold = source_run.font.bold
                                if source_run.font.name:
                                    run.font.name = source_run.font.name
                                if source_run.font.color and source_run.font.color.rgb:
                                    run.font.color.rgb = source_run.font.color.rgb
        
        except Exception as e:
            print(f"⚠️  複製形狀時發生錯誤：{e}")
    
    def replace_text_in_slide(self, slide, replacements):
        """替換投影片中的文字"""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for old_text, new_text in replacements.items():
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
    
    def is_verse_format(self, text):
        """判斷是否為經文格式：〈章節〉 + 經文內容"""
        import re
        pattern = r'^[〈<]([^〉>]+)[〉>]\s*(.+)$'
        match = re.match(pattern, text, re.DOTALL)
        return match
    
    def convert_verse_reference(self, ref):
        """轉換經文章節格式：〈創 19:17〉 → 【創19:17】"""
        # 移除空格
        ref = ref.replace(' ', '').replace('　', '')
        # 轉換括號
        ref = ref.replace('〈', '【').replace('〉', '】')
        ref = ref.replace('<', '【').replace('>', '】')
        return ref
    
    def create_content_slide(self, text):
        """建立內容投影片（使用模板中的內容頁作為參考）"""
        # 選擇適當的模板頁
        # 如果模板有8頁或以上，使用第8頁；否則使用第3頁（簡化版）
        template_index = 7 if len(self.template.slides) > 7 else 2
        if template_index < len(self.template.slides):
            source_slide = self.template.slides[template_index]
            slide_layout = source_slide.slide_layout
        else:
            # 如果連第3頁都沒有，使用第一個可用的版面配置
            slide_layout = self.template.slide_layouts[0]
            source_slide = None
        
        new_slide = self.new_prs.slides.add_slide(slide_layout)
        
        # 刪除所有從模板繼承的文字框（避免空白文字框殘留）
        shapes_to_remove = []
        for shape in new_slide.shapes:
            if hasattr(shape, "text_frame"):
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # 檢查是否為經文格式
        verse_match = self.is_verse_format(text)
        
        # 找到第一個文字框的位置和大小資訊
        source_shape = None
        if source_slide:
            for shape in source_slide.shapes:
                if hasattr(shape, "text_frame"):
                    source_shape = shape
                    break
        
        if source_shape:
            # 調整文字框位置，確保在版面內
            # 使用安全的邊距：左右各 0.5 英吋，上下各 0.3 英吋
            safe_left = Inches(0.5)
            safe_top = Inches(0.3)
            safe_width = Inches(9.0)  # 10 - 0.5*2 = 9
            safe_height = Inches(5.0)  # 5.625 - 0.3*2 ≈ 5
            
            # 建立新的文字框（使用安全範圍）
            new_shape = new_slide.shapes.add_textbox(
                safe_left,
                safe_top,
                safe_width,
                safe_height
            )
            
            # 清空預設文字
            new_shape.text_frame.clear()
            
            # 設定文字框屬性
            new_shape.text_frame.word_wrap = True  # 自動換行
            new_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # 垂直居中對齊
            new_shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE  # 不自動調整大小（避免超出螢幕）
            
            if verse_match:
                # 經文格式：兩個段落，不同顏色
                verse_ref = verse_match.group(1)
                verse_text = verse_match.group(2).strip()
                
                # 轉換章節格式
                verse_ref_formatted = self.convert_verse_reference(verse_ref)
                
                # 第一段：經文章節（淺藍色）
                p1 = new_shape.text_frame.paragraphs[0]
                p1.text = verse_ref_formatted
                for run in p1.runs:
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.name = "微軟正黑體"
                    run.font.color.rgb = RGBColor(121, 155, 193)  # 淺藍色
                
                # 第二段：經文內容（深藍色）
                p2 = new_shape.text_frame.add_paragraph()
                p2.text = verse_text
                for run in p2.runs:
                    run.font.size = Pt(30)
                    run.font.bold = True
                    run.font.name = "微軟正黑體"
                    run.font.color.rgb = RGBColor(27, 54, 106)  # 深藍色
            
            else:
                # 一般文字：單一段落
                p = new_shape.text_frame.paragraphs[0]
                p.text = text
                
                # 複製格式（從模板）
                if source_shape.text_frame.paragraphs:
                    source_p = source_shape.text_frame.paragraphs[0]
                    p.alignment = source_p.alignment
                    
                    # 複製字體格式
                    if source_p.runs:
                        source_run = source_p.runs[0]
                        for target_run in p.runs:
                            if source_run.font.size:
                                target_run.font.size = source_run.font.size
                            if source_run.font.bold:
                                target_run.font.bold = source_run.font.bold
                            if source_run.font.name:
                                target_run.font.name = source_run.font.name
                            if source_run.font.color and source_run.font.color.rgb:
                                target_run.font.color.rgb = source_run.font.color.rgb
        
        return new_slide
    
    def generate(self, output_path):
        """生成完整 PPT（適配不同模板）"""
        print("\n" + "=" * 70)
        print("🎨 開始生成 PPT")
        print("=" * 70)
        
        template_slides_count = len(self.template.slides)
        print(f"📄 模板頁數: {template_slides_count}")
        
        # 準備替換字典
        replacements = {
            '2025年12月31日': self.config.get('DATE', '2025年12月31日'),
            '週三禮拜': self.config.get('SERVICE_TYPE', '週三禮拜'),
            '要避開才能活 ': self.config.get('TITLE', '要避開才能活 '),
            '要避開才能活  這就是天的法則': self.config.get('TITLE', '要避開才能活 這就是天的法則'),
            '這就是天的法則': '這就是天的法則',  # 保留副標題
            '【箴言27章12節、詩篇46篇1節】': f"【{self.config.get('VERSE_REFS', '箴言27章12節、詩篇46篇1節')}】",
        }
        
        if template_slides_count == 4:
            # 簡化版模板（4頁）：封面、主題頁、範例內容頁x2
            print("\n📝 使用簡化版模板（4頁）")
            
            # 1. 複製第1頁（封面）
            print("\n📄 建立封面...")
            slide = self.copy_slide(0)
            self.replace_text_in_slide(slide, replacements)
            
            # 2. 複製第2頁（主題頁）
            print("📄 建立主題頁...")
            slide = self.copy_slide(1)
            self.replace_text_in_slide(slide, replacements)
            
            # 3. 使用第3頁作為內容頁模板，為每段藍色文字創建頁面
            print(f"\n📝 建立內容頁面（{len(self.blue_texts)} 頁）...")
            for i, text in enumerate(self.blue_texts, 1):
                print(f"   建立第 {2+i} 頁：{text[:30]}...")
                self.create_content_slide(text)
            
            # 4. 複製最後一頁（結束頁，如果需要）
            print("\n📄 建立結束頁...")
            slide = self.copy_slide(1)  # 複製主題頁作為結束頁
            self.replace_text_in_slide(slide, replacements)
        
        else:
            # 完整版模板（29頁）
            print("\n📝 使用完整版模板（29頁）")
            
            # 1. 複製並修改前7張固定頁面
            print("\n📄 建立固定頁面（第 1-7 頁）...")
            for i in range(min(7, template_slides_count)):
                print(f"   複製第 {i+1} 頁...")
                slide = self.copy_slide(i)
                self.replace_text_in_slide(slide, replacements)
                
                # 特殊處理第5、6頁（經文內容）
                if i == 4:  # 第5頁
                    verse_ref = self.config.get('VERSE_REF_1', '')
                    verse_text = self.config.get('VERSE_TEXT_1', '')
                    if verse_ref and verse_text:
                        self.replace_text_in_slide(slide, {
                            '【箴言27章12節】': f'【{verse_ref}】',
                            '通達人見禍藏躲；愚蒙人前往受害。': verse_text
                        })
                
                elif i == 5:  # 第6頁
                    verse_ref = self.config.get('VERSE_REF_2', '')
                    verse_text = self.config.get('VERSE_TEXT_2', '')
                    if verse_ref and verse_text:
                        self.replace_text_in_slide(slide, {
                            '【詩篇46篇1節】': f'【{verse_ref}】',
                            ' 神是我們的避難所，是我們的力量，是我們': verse_text
                        })
            
            # 2. 建立藍色文字內容頁（第8頁開始）
            print(f"\n📝 建立內容頁面（第 8-{7+len(self.blue_texts)} 頁）...")
            for i, text in enumerate(self.blue_texts, 1):
                print(f"   建立第 {7+i} 頁：{text[:30]}...")
                self.create_content_slide(text)
            
            # 3. 複製最後2張固定頁面
            print("\n📄 建立結束頁面（最後 2 頁）...")
            for i in [-2, -1]:
                slide_num = len(self.template.slides) + i
                print(f"   複製第 {slide_num+1} 頁...")
                slide = self.copy_slide(slide_num)
                self.replace_text_in_slide(slide, replacements)
        
        # 4. 儲存
        print(f"\n💾 儲存 PPT...")
        try:
            self.new_prs.save(output_path)
            print(f"✅ 成功建立：{output_path}")
            print(f"📊 總共 {len(self.new_prs.slides)} 張投影片")
            return True
        
        except Exception as e:
            print(f"❌ 儲存錯誤：{e}")
            return False


def main():
    """主程式"""
    if len(sys.argv) < 3:
        print("使用方式：")
        print("  新格式（含變數）：")
        print("    python generate_ppt_from_template.py <模板.pptx> <含變數的TXT> [輸出.pptx]")
        print()
        print("  舊格式（分離）：")
        print("    python generate_ppt_from_template.py <模板.pptx> <藍色文字.txt> <配置檔.txt> [輸出.pptx]")
        print()
        print("範例：")
        print("  新格式：")
        print('    python generate_ppt_from_template.py template.pptx 20251231_blue_text.txt output.pptx')
        print()
        print("  舊格式：")
        print('    python generate_ppt_from_template.py "20251231 Wed.pptx" blue_text.txt config.txt output.pptx')
        sys.exit(1)
    
    template_file = sys.argv[1]
    
    # 判斷是新格式還是舊格式
    if len(sys.argv) >= 4 and not sys.argv[3].endswith('.pptx'):
        # 舊格式：4個參數（模板、藍色文字、配置檔、輸出）
        blue_text_file = sys.argv[2]
        config_file = sys.argv[3]
        output_file = sys.argv[4] if len(sys.argv) >= 5 else "output.pptx"
        
        # 檢查檔案
        for file_path in [template_file, blue_text_file, config_file]:
            if not os.path.exists(file_path):
                print(f"❌ 找不到檔案：{file_path}")
                sys.exit(1)
        
        # 生成 PPT（舊格式）
        generator = PPTGenerator(template_file)
        
        if not generator.load_config(config_file):
            sys.exit(1)
        
        if not generator.load_blue_texts(blue_text_file):
            sys.exit(1)
    else:
        # 新格式：3個參數（模板、含變數的TXT、輸出）
        blue_text_file = sys.argv[2]
        output_file = sys.argv[3] if len(sys.argv) >= 4 else "output.pptx"
        
        # 檢查檔案
        for file_path in [template_file, blue_text_file]:
            if not os.path.exists(file_path):
                print(f"❌ 找不到檔案：{file_path}")
                sys.exit(1)
        
        # 生成 PPT（新格式：從TXT中讀取變數和內容）
        generator = PPTGenerator(template_file)
        
        if not generator.load_blue_texts(blue_text_file):
            sys.exit(1)
    
    if generator.generate(output_file):
        print("\n" + "=" * 70)
        print("🎉 完成！")
        print("=" * 70)
    else:
        sys.exit(1)


if __name__ == "__main__":
    main()
