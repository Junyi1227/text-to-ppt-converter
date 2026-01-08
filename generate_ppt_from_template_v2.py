#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT ç”Ÿæˆç¨‹å¼ V2 - åŸºæ–¼ template.pptx çš„å½ˆæ€§åŒ–ç‰ˆæœ¬

ä½¿ç”¨æ–¹å¼ï¼š
    python generate_ppt_from_template_v2.py template.pptx input.txt config.txt output.pptx

åŠŸèƒ½ï¼š
    - æ”¯æ´å½ˆæ€§åŒ–çš„é é¢çµæ§‹å®šç¾©ï¼ˆé€é configï¼‰
    - æ”¯æ´è®Šæ•¸æ¨¡æ¿ï¼ˆå¾ TXT è®€å–ï¼‰
    - è‡ªå‹•è­˜åˆ¥ç¶“æ–‡æ ¼å¼
    - æ”¯æ´å¤šç¨®é é¢é¡å‹ï¼šCOVER, TITLE, CONTENT, BIBLE, AUTOCONTENT
"""

import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


class PPTGeneratorV2:
    """PPT ç”Ÿæˆå™¨ V2"""
    
    def __init__(self, template_path, output_path):
        """
        åˆå§‹åŒ– PPT ç”Ÿæˆå™¨
        
        Args:
            template_path: æ¨¡æ¿ PPT è·¯å¾‘ï¼ˆå¿…é ˆåŒ…å« 4 é ï¼‰
            output_path: è¼¸å‡º PPT è·¯å¾‘
        """
        # å…ˆè¤‡è£½ template åˆ° output
        import shutil
        shutil.copy2(template_path, output_path)
        
        # é–‹å•Ÿè¼¸å‡ºæª”æ¡ˆï¼ˆåŒ…å«æ¨¡æ¿çš„ 4 é ï¼‰
        self.output_prs = Presentation(output_path)
        self.output_path = output_path
        
        # ç¢ºèªæ¨¡æ¿æœ‰ 4 é 
        if len(self.output_prs.slides) < 4:
            raise ValueError(f"æ¨¡æ¿å¿…é ˆåŒ…å«è‡³å°‘ 4 é ï¼Œç›®å‰åªæœ‰ {len(self.output_prs.slides)} é ")
        
        # æ³¨æ„ï¼šä¸åˆªé™¤æ¨¡æ¿é ï¼Œç¨å¾Œç”Ÿæˆæ™‚æœƒç”¨åˆ°
        
        # è®Šæ•¸å­—å…¸
        self.variables = {}
        # å…§å®¹åˆ—è¡¨
        self.content_lines = []
        # é é¢çµæ§‹
        self.page_structure = []
        # è¨˜éŒ„éœ€è¦åˆªé™¤çš„æ¨¡æ¿é ç´¢å¼•
        self.template_page_count = len(self.output_prs.slides)
    
    def load_variables_and_content(self, txt_path):
        """
        å¾ TXT æª”æ¡ˆè®€å–è®Šæ•¸å’Œå…§å®¹ï¼ˆä½¿ç”¨ç©ºè¡Œåˆ†éš”é é¢ï¼‰
        
        Args:
            txt_path: TXT æª”æ¡ˆè·¯å¾‘
        """
        with open(txt_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        in_variables = False
        in_content = False
        current_block = []
        
        for line in lines:
            line = line.rstrip('\n')
            
            # æª¢æŸ¥è®Šæ•¸å€é–‹å§‹
            if line.strip() == '[è®Šæ•¸]':
                in_variables = True
                continue
            
            # æª¢æŸ¥è®Šæ•¸å€çµæŸ
            if line.strip() == '[è®Šæ•¸çµæŸ]':
                in_variables = False
                in_content = True
                continue
            
            # è®€å–è®Šæ•¸
            if in_variables and '=' in line:
                key, value = line.split('=', 1)
                self.variables[key.strip()] = value.strip()
            
            # è®€å–å…§å®¹ï¼ˆä½¿ç”¨ç©ºè¡Œåˆ†éš”ä¸åŒé é¢ï¼‰
            elif in_content:
                if line.strip():
                    # æœ‰å…§å®¹çš„è¡Œï¼ŒåŠ å…¥ç•¶å‰å€å¡Š
                    current_block.append(line.strip())
                else:
                    # ç©ºè¡Œï¼Œè¡¨ç¤ºä¸€å€‹å€å¡ŠçµæŸ
                    if current_block:
                        # å°‡å€å¡Šåˆä½µæˆä¸€å€‹é …ç›®ï¼ˆç”¨æ›è¡Œç¬¦é€£æ¥ï¼‰
                        self.content_lines.append('\n'.join(current_block))
                        current_block = []
        
        # è™•ç†æœ€å¾Œä¸€å€‹å€å¡Šï¼ˆå¦‚æœæª”æ¡ˆçµå°¾æ²’æœ‰ç©ºè¡Œï¼‰
        if current_block:
            self.content_lines.append('\n'.join(current_block))
        
        print(f"âœ… è®€å–è®Šæ•¸: {len(self.variables)} å€‹")
        print(f"âœ… è®€å–å…§å®¹å€å¡Š: {len(self.content_lines)} å€‹ï¼ˆç”¨ç©ºè¡Œåˆ†éš”ï¼‰")
    
    def load_config(self, config_path):
        """
        å¾ config æª”æ¡ˆè®€å–é é¢çµæ§‹
        
        Args:
            config_path: config æª”æ¡ˆè·¯å¾‘
        """
        with open(config_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        in_structure = False
        
        for line in lines:
            line = line.strip()
            
            # è·³éç©ºè¡Œå’Œè¨»è§£
            if not line or line.startswith('#'):
                continue
            
            # æª¢æŸ¥é é¢çµæ§‹å€é–‹å§‹
            if line == '[é é¢çµæ§‹]':
                in_structure = True
                continue
            
            # è®€å–é é¢çµæ§‹
            if in_structure:
                # è§£æé é¢é¡å‹å’Œåƒæ•¸
                if '=' in line:
                    parts = line.split('=', 1)
                    page_type = parts[0].strip()
                    param = parts[1].strip()
                    self.page_structure.append((page_type, param))
                else:
                    page_type = line.strip()
                    self.page_structure.append((page_type, None))
        
        print(f"âœ… è®€å–é é¢çµæ§‹: {len(self.page_structure)} é ")
    
    def is_verse_format(self, text):
        """
        åˆ¤æ–·æ˜¯å¦ç‚ºç¶“æ–‡æ ¼å¼
        æ”¯æ´å…©ç¨®æ ¼å¼ï¼š
        1. å–®è¡Œï¼šã€ˆå‰µ19:17ã€‰é ˜ä»–å€‘å‡ºä¾†...
        2. å¤šè¡Œï¼šç¬¬ä¸€è¡Œæ˜¯ç« ç¯€ï¼Œç¬¬äºŒè¡Œæ˜¯å…§å®¹
        
        Args:
            text: è¦åˆ¤æ–·çš„æ–‡å­—
            
        Returns:
            Match object å¦‚æœåŒ¹é…ï¼Œå¦å‰‡ None
        """
        # å–®è¡Œæ ¼å¼ï¼šã€ˆç« ç¯€ã€‰å…§å®¹
        pattern = r'^[ã€ˆ<]([^ã€‰>]+)[ã€‰>](.+)$'
        match = re.match(pattern, text)
        return match
    
    def convert_verse_reference(self, verse_ref):
        """
        è½‰æ›ç¶“æ–‡ç« ç¯€æ ¼å¼
        å‰µ19:17 â†’ å‰µä¸–è¨˜19ç« 17ç¯€
        ç®´è¨€27ç« 12ç¯€ â†’ ç®´è¨€27ç« 12ç¯€ï¼ˆä¸è®Šï¼‰
        
        Args:
            verse_ref: åŸå§‹ç« ç¯€æ ¼å¼
            
        Returns:
            è½‰æ›å¾Œçš„ç« ç¯€æ ¼å¼
        """
        # å¦‚æœå·²ç¶“åŒ…å«ã€Œç« ã€ã€Œç¯€ã€ï¼Œç›´æ¥è¿”å›
        if 'ç« ' in verse_ref and 'ç¯€' in verse_ref:
            return verse_ref
        
        # è½‰æ›ç°¡åŒ–æ ¼å¼ï¼ˆä¾‹å¦‚ï¼šå‰µ19:17ï¼‰
        pattern = r'^([^0-9]+)(\d+):(\d+)$'
        match = re.match(pattern, verse_ref)
        
        if match:
            book = match.group(1)
            chapter = match.group(2)
            verse = match.group(3)
            
            # æ›¸å·åç¨±è½‰æ›ï¼ˆå¦‚æœéœ€è¦ï¼‰
            book_map = {
                'å‰µ': 'å‰µä¸–è¨˜',
                'å‡º': 'å‡ºåŸƒåŠè¨˜',
                'åˆ©': 'åˆ©æœªè¨˜',
                'æ°‘': 'æ°‘æ•¸è¨˜',
                'ç”³': 'ç”³å‘½è¨˜',
                'æ›¸': 'ç´„æ›¸äºè¨˜',
                'å£«': 'å£«å¸«è¨˜',
                'å¾—': 'è·¯å¾—è¨˜',
                'æ’’ä¸Š': 'æ’’æ¯è€³è¨˜ä¸Š',
                'æ’’ä¸‹': 'æ’’æ¯è€³è¨˜ä¸‹',
                'ç‹ä¸Š': 'åˆ—ç‹ç´€ä¸Š',
                'ç‹ä¸‹': 'åˆ—ç‹ç´€ä¸‹',
                'ä»£ä¸Š': 'æ­·ä»£å¿—ä¸Š',
                'ä»£ä¸‹': 'æ­·ä»£å¿—ä¸‹',
                'æ‹‰': 'ä»¥æ–¯æ‹‰è¨˜',
                'å°¼': 'å°¼å¸Œç±³è¨˜',
                'æ–¯': 'ä»¥æ–¯å¸–è¨˜',
                'ä¼¯': 'ç´„ä¼¯è¨˜',
                'è©©': 'è©©ç¯‡',
                'ç®´': 'ç®´è¨€',
                'å‚³': 'å‚³é“æ›¸',
                'æ­Œ': 'é›…æ­Œ',
                'è³½': 'ä»¥è³½äºæ›¸',
                'è€¶': 'è€¶åˆ©ç±³æ›¸',
                'å“€': 'è€¶åˆ©ç±³å“€æ­Œ',
                'çµ': 'ä»¥è¥¿çµæ›¸',
                'ä½†': 'ä½†ä»¥ç†æ›¸',
                'ä½•': 'ä½•è¥¿é˜¿æ›¸',
                'ç¥': 'ç´„ç¥æ›¸',
                'æ‘©': 'é˜¿æ‘©å¸æ›¸',
                'ä¿„': 'ä¿„å·´åº•äºæ›¸',
                'æ‹¿': 'ç´„æ‹¿æ›¸',
                'å½Œ': 'å½Œè¿¦æ›¸',
                'é´»': 'é‚£é´»æ›¸',
                'å“ˆ': 'å“ˆå·´è°·æ›¸',
                'ç•ª': 'è¥¿ç•ªé›…æ›¸',
                'è©²': 'å“ˆè©²æ›¸',
                'äº': 'æ’’è¿¦åˆ©äºæ›¸',
                'ç‘ª': 'ç‘ªæ‹‰åŸºæ›¸',
                'å¤ª': 'é¦¬å¤ªç¦éŸ³',
                'å¯': 'é¦¬å¯ç¦éŸ³',
                'è·¯': 'è·¯åŠ ç¦éŸ³',
                'ç´„': 'ç´„ç¿°ç¦éŸ³',
                'å¾’': 'ä½¿å¾’è¡Œå‚³',
                'ç¾…': 'ç¾…é¦¬æ›¸',
                'æ—å‰': 'å“¥æ—å¤šå‰æ›¸',
                'æ—å¾Œ': 'å“¥æ—å¤šå¾Œæ›¸',
                'åŠ ': 'åŠ æ‹‰å¤ªæ›¸',
                'å¼—': 'ä»¥å¼—æ‰€æ›¸',
                'è…“': 'è…“ç«‹æ¯”æ›¸',
                'è¥¿': 'æ­Œç¾…è¥¿æ›¸',
                'å¸–å‰': 'å¸–æ’’ç¾…å°¼è¿¦å‰æ›¸',
                'å¸–å¾Œ': 'å¸–æ’’ç¾…å°¼è¿¦å¾Œæ›¸',
                'æå‰': 'ææ‘©å¤ªå‰æ›¸',
                'æå¾Œ': 'ææ‘©å¤ªå¾Œæ›¸',
                'å¤š': 'æå¤šæ›¸',
                'é–€': 'è…“åˆ©é–€æ›¸',
                'ä¾†': 'å¸Œä¼¯ä¾†æ›¸',
                'é›…': 'é›…å„æ›¸',
                'å½¼å‰': 'å½¼å¾—å‰æ›¸',
                'å½¼å¾Œ': 'å½¼å¾—å¾Œæ›¸',
                'ç´„å£¹': 'ç´„ç¿°ä¸€æ›¸',
                'ç´„è²³': 'ç´„ç¿°äºŒæ›¸',
                'ç´„åƒ': 'ç´„ç¿°ä¸‰æ›¸',
                'çŒ¶': 'çŒ¶å¤§æ›¸',
                'å•Ÿ': 'å•Ÿç¤ºéŒ„'
            }
            
            full_book = book_map.get(book, book)
            return f"{full_book}{chapter}ç« {verse}ç¯€"
        
        return verse_ref
    
    def create_cover_page(self, subtitle=None):
        """
        å»ºç«‹å°é¢é ï¼ˆè¤‡è£½ template ç¬¬ 1 é ä¸¦ä¿®æ”¹å…§å®¹ï¼‰
        
        Args:
            subtitle: å°æ¨™é¡Œï¼ˆå¯é¸ï¼‰
        """
        # ä½¿ç”¨æ¨¡æ¿ç¬¬ 1 é çš„ç‰ˆé¢é…ç½®
        template_slide = self.output_prs.slides[0]
        slide_layout = template_slide.slide_layout
        new_slide = self.output_prs.slides.add_slide(slide_layout)
        
        # åˆªé™¤å¾ç‰ˆé¢é…ç½®ç¹¼æ‰¿çš„ç©ºæ–‡å­—æ¡†
        shapes_to_remove = []
        for shape in new_slide.shapes:
            if hasattr(shape, "text_frame") and not shape.text.strip():
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # è¤‡è£½æ¨¡æ¿é çš„æ‰€æœ‰å½¢ç‹€ä¸¦ä¿®æ”¹æ–‡å­—
        for shape in template_slide.shapes:
            if hasattr(shape, "text_frame"):
                # æ ¹æ“šä½ç½®åˆ¤æ–·æ˜¯å“ªå€‹æ–‡å­—æ¡†
                # æ–‡å­—æ¡†1: æ—¥æœŸ+ç¦®æ‹œé¡å‹ (top â‰ˆ 1.23")
                # æ–‡å­—æ¡†2: å°æ¨™é¡Œ (top â‰ˆ 4.30")
                # æ–‡å­—æ¡†3: ç¶“æ–‡ç« ç¯€ (top â‰ˆ 3.40")
                
                if abs(shape.top.inches - 1.23) < 0.1:
                    # æ–‡å­—æ¡†1: æ—¥æœŸ+ç¦®æ‹œé¡å‹
                    date = self.variables.get('æ—¥æœŸ', '')
                    service_type = self.variables.get('ç¦®æ‹œé¡å‹', '')
                    text = f"{date}\n\n{service_type}"
                    self._create_textbox_with_format(new_slide, shape, text)
                
                elif abs(shape.top.inches - 4.30) < 0.1:
                    # æ–‡å­—æ¡†2: å°æ¨™é¡Œï¼ˆåªæœ‰åœ¨æœ‰åƒæ•¸æ™‚æ‰é¡¯ç¤ºï¼‰
                    if subtitle:
                        self._create_textbox_with_format(new_slide, shape, subtitle)
                
                elif abs(shape.top.inches - 3.40) < 0.1:
                    # æ–‡å­—æ¡†3: ç¶“æ–‡ç« ç¯€
                    verse_refs = self.variables.get('ç¶“æ–‡ç« ç¯€', '')
                    self._create_textbox_with_format(new_slide, shape, verse_refs)
        
        return new_slide
    
    def create_title_page(self, subtitle=None):
        """
        å»ºç«‹ä¸»é¡Œé ï¼ˆè¤‡è£½ template ç¬¬ 2 é ä¸¦ä¿®æ”¹å…§å®¹ï¼‰
        
        Args:
            subtitle: å°æ¨™é¡Œï¼ˆå¯é¸ï¼‰
        """
        # ä½¿ç”¨æ¨¡æ¿ç¬¬ 2 é çš„ç‰ˆé¢é…ç½®
        template_slide = self.output_prs.slides[1]
        slide_layout = template_slide.slide_layout
        new_slide = self.output_prs.slides.add_slide(slide_layout)
        
        # åˆªé™¤å¾ç‰ˆé¢é…ç½®ç¹¼æ‰¿çš„ç©ºæ–‡å­—æ¡†
        shapes_to_remove = []
        for shape in new_slide.shapes:
            if hasattr(shape, "text_frame") and not shape.text.strip():
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # è¤‡è£½æ¨¡æ¿é çš„æ‰€æœ‰å½¢ç‹€ä¸¦ä¿®æ”¹æ–‡å­—
        for shape in template_slide.shapes:
            if hasattr(shape, "text_frame"):
                # æ–‡å­—æ¡†1: æ—¥æœŸ+ç¦®æ‹œé¡å‹ (top â‰ˆ 0.51")
                # æ–‡å­—æ¡†2: ä¸»é¡Œ (top â‰ˆ 1.72")
                # æ–‡å­—æ¡†3: ç¶“æ–‡ç« ç¯€ (top â‰ˆ 3.76")
                # æ–‡å­—æ¡†4: å°æ¨™é¡Œ (top â‰ˆ 4.46")
                
                if abs(shape.top.inches - 0.51) < 0.1:
                    # æ–‡å­—æ¡†1: æ—¥æœŸ+ç¦®æ‹œé¡å‹
                    date = self.variables.get('æ—¥æœŸ', '')
                    service_type = self.variables.get('ç¦®æ‹œé¡å‹', '')
                    text = f"{date} {service_type}"
                    self._create_textbox_with_format(new_slide, shape, text)
                
                elif abs(shape.top.inches - 1.72) < 0.1:
                    # æ–‡å­—æ¡†2: ä¸»é¡Œ
                    title = self.variables.get('ä¸»é¡Œ', '')
                    self._create_textbox_with_format(new_slide, shape, title)
                
                elif abs(shape.top.inches - 3.76) < 0.1:
                    # æ–‡å­—æ¡†3: ç¶“æ–‡ç« ç¯€
                    verse_refs = self.variables.get('ç¶“æ–‡ç« ç¯€', '')
                    self._create_textbox_with_format(new_slide, shape, verse_refs)
                
                elif abs(shape.top.inches - 4.46) < 0.1:
                    # æ–‡å­—æ¡†4: å°æ¨™é¡Œï¼ˆåªæœ‰åœ¨æœ‰åƒæ•¸æ™‚æ‰é¡¯ç¤ºï¼‰
                    if subtitle:
                        self._create_textbox_with_format(new_slide, shape, subtitle)
        
        return new_slide
    
    def create_content_page(self, text):
        """
        å»ºç«‹å…§æ–‡é ï¼ˆè¤‡è£½ template ç¬¬ 3 é ä¸¦ä¿®æ”¹å…§å®¹ï¼‰
        
        Args:
            text: å…§å®¹æ–‡å­—
        """
        # ä½¿ç”¨æ¨¡æ¿ç¬¬ 3 é çš„ç‰ˆé¢é…ç½®
        template_slide = self.output_prs.slides[2]
        slide_layout = template_slide.slide_layout
        new_slide = self.output_prs.slides.add_slide(slide_layout)
        
        # åˆªé™¤å¾ç‰ˆé¢é…ç½®ç¹¼æ‰¿çš„ç©ºæ–‡å­—æ¡†
        shapes_to_remove = []
        for shape in new_slide.shapes:
            if hasattr(shape, "text_frame") and not shape.text.strip():
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # æ‰¾åˆ°æ¨¡æ¿é çš„ç¬¬ä¸€å€‹æ–‡å­—æ¡†ä¸¦è¤‡è£½
        for shape in template_slide.shapes:
            if hasattr(shape, "text_frame"):
                # ä½¿ç”¨æ¨¡æ¿çš„ä½ç½®å’Œå¤§å°ï¼ˆä¸è¦å¯«æ­»ï¼‰
                self._create_textbox_with_format(new_slide, shape, text)
                break
        
        return new_slide
    
    def create_verse_page(self, verse_ref, verse_text):
        """
        å»ºç«‹ç¶“æ–‡é ï¼ˆè¤‡è£½ template ç¬¬ 4 é ä¸¦ä¿®æ”¹å…§å®¹ï¼‰
        
        Args:
            verse_ref: ç¶“æ–‡ç« ç¯€
            verse_text: ç¶“æ–‡å…§å®¹
        """
        # ä½¿ç”¨æ¨¡æ¿ç¬¬ 4 é çš„ç‰ˆé¢é…ç½®
        template_slide = self.output_prs.slides[3]
        slide_layout = template_slide.slide_layout
        new_slide = self.output_prs.slides.add_slide(slide_layout)
        
        # åˆªé™¤å¾ç‰ˆé¢é…ç½®ç¹¼æ‰¿çš„ç©ºæ–‡å­—æ¡†
        shapes_to_remove = []
        for shape in new_slide.shapes:
            if hasattr(shape, "text_frame") and not shape.text.strip():
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        # æ‰¾åˆ°æ¨¡æ¿é çš„ç¬¬ä¸€å€‹æ–‡å­—æ¡†
        source_shape = None
        for shape in template_slide.shapes:
            if hasattr(shape, "text_frame"):
                source_shape = shape
                break
        
        if source_shape:
            # ä½¿ç”¨æ¨¡æ¿çš„ä½ç½®å’Œå¤§å°ï¼ˆä¸è¦å¯«æ­»ï¼‰
            new_shape = new_slide.shapes.add_textbox(
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )
            
            # æ¸…ç©ºé è¨­æ–‡å­—
            new_shape.text_frame.clear()
            
            # è¤‡è£½æ–‡å­—æ¡†å±¬æ€§
            new_shape.text_frame.word_wrap = source_shape.text_frame.word_wrap
            new_shape.text_frame.vertical_anchor = source_shape.text_frame.vertical_anchor
            new_shape.text_frame.auto_size = source_shape.text_frame.auto_size
            
            # è½‰æ›ç« ç¯€æ ¼å¼
            verse_ref_formatted = self.convert_verse_reference(verse_ref)
            
            # ç¬¬ä¸€æ®µï¼šç¶“æ–‡ç« ç¯€ï¼ˆå¾æ¨¡æ¿è¤‡è£½æ ¼å¼ï¼‰
            p1 = new_shape.text_frame.paragraphs[0]
            p1.text = verse_ref_formatted
            
            # è¤‡è£½ç¬¬ä¸€æ®µæ ¼å¼ï¼ˆå¦‚æœæ¨¡æ¿æœ‰çš„è©±ï¼‰
            if source_shape.text_frame.paragraphs:
                source_p = source_shape.text_frame.paragraphs[0]
                p1.alignment = source_p.alignment
                
                if source_p.runs:
                    source_run = source_p.runs[0]
                    for run in p1.runs:
                        if source_run.font.size:
                            run.font.size = source_run.font.size
                        if source_run.font.bold is not None:
                            run.font.bold = source_run.font.bold
                        if source_run.font.name:
                            run.font.name = source_run.font.name
                        # ç¶“æ–‡ç« ç¯€ä½¿ç”¨æ·ºè—è‰²
                        run.font.color.rgb = RGBColor(121, 155, 193)
            
            # ç¬¬äºŒæ®µï¼šç¶“æ–‡å…§å®¹
            p2 = new_shape.text_frame.add_paragraph()
            p2.text = verse_text
            
            # è¤‡è£½ç¬¬äºŒæ®µæ ¼å¼ï¼ˆå¦‚æœæ¨¡æ¿æœ‰å¤šå€‹æ®µè½çš„è©±ï¼‰
            if len(source_shape.text_frame.paragraphs) > 1:
                source_p2 = source_shape.text_frame.paragraphs[1]
                p2.alignment = source_p2.alignment
                
                if source_p2.runs:
                    source_run2 = source_p2.runs[0]
                    for run in p2.runs:
                        if source_run2.font.size:
                            run.font.size = source_run2.font.size
                        if source_run2.font.bold is not None:
                            run.font.bold = source_run2.font.bold
                        if source_run2.font.name:
                            run.font.name = source_run2.font.name
                        # ç¶“æ–‡å…§å®¹ä½¿ç”¨æ·±è—è‰²
                        run.font.color.rgb = RGBColor(27, 54, 106)
            else:
                # å¦‚æœæ¨¡æ¿åªæœ‰ä¸€æ®µï¼Œä½¿ç”¨ç¬¬ä¸€æ®µçš„æ ¼å¼
                if source_shape.text_frame.paragraphs:
                    source_p = source_shape.text_frame.paragraphs[0]
                    p2.alignment = source_p.alignment
                    
                    if source_p.runs:
                        source_run = source_p.runs[0]
                        for run in p2.runs:
                            if source_run.font.size:
                                run.font.size = source_run.font.size
                            if source_run.font.bold is not None:
                                run.font.bold = source_run.font.bold
                            if source_run.font.name:
                                run.font.name = source_run.font.name
                            # ç¶“æ–‡å…§å®¹ä½¿ç”¨æ·±è—è‰²
                            run.font.color.rgb = RGBColor(27, 54, 106)
        
        return new_slide
    
    def _create_textbox_with_format(self, slide, source_shape, text):
        """
        å‰µå»ºæ–‡å­—æ¡†ä¸¦è¤‡è£½æ ¼å¼ï¼ˆæ”¯æ´å¤šæ®µè½ï¼‰
        
        Args:
            slide: ç›®æ¨™æŠ•å½±ç‰‡
            source_shape: ä¾†æºå½¢ç‹€ï¼ˆç”¨æ–¼è¤‡è£½ä½ç½®å’Œæ ¼å¼ï¼‰
            text: è¦å¡«å…¥çš„æ–‡å­—
        """
        # å‰µå»ºæ–°æ–‡å­—æ¡†
        new_shape = slide.shapes.add_textbox(
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )
        
        # è¨­å®šæ–‡å­—
        new_shape.text = text
        
        # è¤‡è£½æ–‡å­—æ¡†å±¬æ€§
        new_shape.text_frame.word_wrap = source_shape.text_frame.word_wrap
        new_shape.text_frame.vertical_anchor = source_shape.text_frame.vertical_anchor
        new_shape.text_frame.auto_size = source_shape.text_frame.auto_size
        
        # è¤‡è£½æ‰€æœ‰æ®µè½çš„æ ¼å¼
        text_paragraphs = text.split('\n')
        target_paragraphs = new_shape.text_frame.paragraphs
        
        # ç¢ºä¿ç›®æ¨™æœ‰è¶³å¤ çš„æ®µè½
        while len(target_paragraphs) < len(text_paragraphs):
            new_shape.text_frame.add_paragraph()
            target_paragraphs = new_shape.text_frame.paragraphs
        
        # ç‚ºæ¯å€‹æ®µè½è¤‡è£½å°æ‡‰çš„æ ¼å¼
        for i, target_p in enumerate(target_paragraphs):
            # æ‰¾åˆ°å°æ‡‰çš„æºæ®µè½ï¼ˆå¦‚æœæ²’æœ‰å°±ç”¨æœ€å¾Œä¸€å€‹ï¼‰
            source_para_index = min(i, len(source_shape.text_frame.paragraphs) - 1)
            if source_para_index >= 0 and source_para_index < len(source_shape.text_frame.paragraphs):
                source_p = source_shape.text_frame.paragraphs[source_para_index]
                
                # è¤‡è£½æ®µè½å°é½Š
                target_p.alignment = source_p.alignment
                
                # è¤‡è£½å­—é«”æ ¼å¼
                if source_p.runs and target_p.runs:
                    source_run = source_p.runs[0]
                    for target_run in target_p.runs:
                        if source_run.font.size:
                            target_run.font.size = source_run.font.size
                        if source_run.font.bold is not None:
                            target_run.font.bold = source_run.font.bold
                        if source_run.font.name:
                            target_run.font.name = source_run.font.name
                        if source_run.font.color and source_run.font.color.rgb:
                            target_run.font.color.rgb = source_run.font.color.rgb
        
        return new_shape
    
    def _copy_text_format(self, source_shape, target_shape):
        """
        è¤‡è£½æ–‡å­—æ ¼å¼
        
        Args:
            source_shape: ä¾†æºå½¢ç‹€
            target_shape: ç›®æ¨™å½¢ç‹€
        """
        if not source_shape.text_frame.paragraphs or not target_shape.text_frame.paragraphs:
            return
        
        source_p = source_shape.text_frame.paragraphs[0]
        target_p = target_shape.text_frame.paragraphs[0]
        
        # è¤‡è£½æ®µè½å°é½Š
        target_p.alignment = source_p.alignment
        
        # è¤‡è£½å­—é«”æ ¼å¼
        if source_p.runs and target_p.runs:
            source_run = source_p.runs[0]
            for target_run in target_p.runs:
                if source_run.font.size:
                    target_run.font.size = source_run.font.size
                if source_run.font.bold is not None:
                    target_run.font.bold = source_run.font.bold
                if source_run.font.name:
                    target_run.font.name = source_run.font.name
                if source_run.font.color and source_run.font.color.rgb:
                    target_run.font.color.rgb = source_run.font.color.rgb
    
    def generate(self):
        """
        æ ¹æ“šé é¢çµæ§‹ç”Ÿæˆ PPT
        """
        content_index = 0  # è¿½è¹¤ AUTOCONTENT çš„ç•¶å‰ç´¢å¼•
        
        for page_type, param in self.page_structure:
            print(f"ç”Ÿæˆé é¢: {page_type}" + (f" = {param}" if param else ""))
            
            if page_type == "COVER":
                # å°é¢é 
                self.create_cover_page(subtitle=param)
            
            elif page_type == "TITLE":
                # ä¸»é¡Œé 
                self.create_title_page(subtitle=param)
            
            elif page_type == "CONTENT":
                # å…§æ–‡é ï¼ˆå›ºå®šå…§å®¹ï¼‰
                if param:
                    self.create_content_page(param)
            
            elif page_type == "BIBLE":
                # ç¶“æ–‡é ï¼ˆè®€å–è®Šæ•¸å€çš„ç¶“æ–‡1, ç¶“æ–‡2, ...ï¼‰
                verse_num = 1
                while True:
                    verse_key = f"ç¶“æ–‡{verse_num}"
                    if verse_key not in self.variables:
                        break
                    
                    verse_data = self.variables[verse_key]
                    # ç”¨ ã€‰ åˆ†éš”ç« ç¯€å’Œå…§å®¹
                    if 'ã€‰' in verse_data:
                        verse_ref, verse_text = verse_data.split('ã€‰', 1)
                        verse_ref = verse_ref.lstrip('ã€ˆ<')
                        verse_text = verse_text.strip()
                        
                        print(f"  ç”Ÿæˆç¶“æ–‡é  {verse_num}: {verse_ref}")
                        self.create_verse_page(verse_ref, verse_text)
                    
                    verse_num += 1
            
            elif page_type == "AUTOCONTENT":
                # è‡ªå‹•å…§å®¹é ï¼ˆå¾å…§å®¹å€è®€å–ï¼Œæ¯å€‹å€å¡Šæ˜¯ä¸€é ï¼‰
                while content_index < len(self.content_lines):
                    block = self.content_lines[content_index]
                    content_index += 1
                    
                    # æª¢æŸ¥å€å¡Šçš„ç¬¬ä¸€è¡Œæ˜¯å¦ç‚ºç¶“æ–‡æ ¼å¼
                    lines_in_block = block.split('\n')
                    first_line = lines_in_block[0] if lines_in_block else ""
                    
                    # æª¢æŸ¥æ˜¯å¦ç‚ºç¶“æ–‡æ ¼å¼ï¼ˆå–®è¡Œï¼‰
                    verse_match = self.is_verse_format(first_line)
                    if verse_match:
                        # å–®è¡Œç¶“æ–‡æ ¼å¼ï¼šã€ˆç« ç¯€ã€‰å…§å®¹
                        verse_ref = verse_match.group(1)
                        verse_text = verse_match.group(2).strip()
                        print(f"  ç”Ÿæˆç¶“æ–‡é : {verse_ref}")
                        self.create_verse_page(verse_ref, verse_text)
                    elif first_line.startswith('ã€ˆ') or first_line.startswith('<'):
                        # å¤šè¡Œç¶“æ–‡æ ¼å¼ï¼šç¬¬ä¸€è¡Œæ˜¯ç« ç¯€ï¼Œå¾Œé¢æ˜¯å…§å®¹
                        verse_ref = first_line.lstrip('ã€ˆ<').rstrip('ã€‰>')
                        verse_text = '\n'.join(lines_in_block[1:]) if len(lines_in_block) > 1 else ""
                        print(f"  ç”Ÿæˆç¶“æ–‡é : {verse_ref}")
                        self.create_verse_page(verse_ref, verse_text)
                    else:
                        # ä¸€èˆ¬å…§å®¹ï¼ˆæ•´å€‹å€å¡Šï¼‰
                        print(f"  ç”Ÿæˆå…§æ–‡é ")
                        self.create_content_page(block)
        
        # åˆªé™¤å‰é¢çš„æ¨¡æ¿é ï¼ˆ4 é ï¼‰
        print(f"\nåˆªé™¤æ¨¡æ¿é ...")
        for i in range(self.template_page_count - 1, -1, -1):
            rId = self.output_prs.slides._sldIdLst[i].rId
            self.output_prs.part.drop_rel(rId)
            del self.output_prs.slides._sldIdLst[i]
        
        # å„²å­˜ PPT
        self.output_prs.save(self.output_path)
        print(f"\nâœ… PPT ç”Ÿæˆå®Œæˆï¼")
        print(f"ğŸ“Š ç¸½å…±ç”Ÿæˆ {len(self.output_prs.slides)} å¼µæŠ•å½±ç‰‡")
        print(f"ğŸ’¾ å·²å„²å­˜åˆ°ï¼š{self.output_path}")


def main():
    """ä¸»ç¨‹å¼"""
    if len(sys.argv) != 5:
        print("ä½¿ç”¨æ–¹å¼ï¼š")
        print("  python generate_ppt_from_template_v2.py template.pptx input.txt config.txt output.pptx")
        print()
        print("åƒæ•¸èªªæ˜ï¼š")
        print("  template.pptx  - æ¨¡æ¿ PPTï¼ˆå¿…é ˆåŒ…å« 4 é ï¼‰")
        print("  input.txt      - è¼¸å…¥æ–‡å­—æª”ï¼ˆåŒ…å«è®Šæ•¸å’Œå…§å®¹ï¼‰")
        print("  config.txt     - è¨­å®šæª”ï¼ˆå®šç¾©é é¢çµæ§‹ï¼‰")
        print("  output.pptx    - è¼¸å‡º PPT æª”å")
        sys.exit(1)
    
    template_path = sys.argv[1]
    input_path = sys.argv[2]
    config_path = sys.argv[3]
    output_path = sys.argv[4]
    
    print("=" * 60)
    print("PPT ç”Ÿæˆç¨‹å¼ V2")
    print("=" * 60)
    print(f"æ¨¡æ¿æª”æ¡ˆï¼š{template_path}")
    print(f"è¼¸å…¥æ–‡å­—ï¼š{input_path}")
    print(f"è¨­å®šæª”æ¡ˆï¼š{config_path}")
    print(f"è¼¸å‡ºæª”æ¡ˆï¼š{output_path}")
    print("=" * 60)
    print()
    
    try:
        # å»ºç«‹ç”Ÿæˆå™¨ï¼ˆæœƒå…ˆè¤‡è£½ template åˆ° outputï¼‰
        generator = PPTGeneratorV2(template_path, output_path)
        
        # è¼‰å…¥è®Šæ•¸å’Œå…§å®¹
        generator.load_variables_and_content(input_path)
        
        # è¼‰å…¥è¨­å®š
        generator.load_config(config_path)
        
        # ç”Ÿæˆ PPT
        generator.generate()
        
    except Exception as e:
        print(f"âŒ éŒ¯èª¤ï¼š{e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
