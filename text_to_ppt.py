#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文字轉 PowerPoint 轉換器 (跨平台版本)
Text to PowerPoint Converter (Cross-Platform)

支援 Mac / Windows / Linux
需要安裝：pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import sys
import os


class TextToPPTConverter:
    """文字轉 PPT 轉換器"""
    
    def __init__(self):
        self.prs = Presentation()
        # 設定投影片尺寸 (16:9)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
        # 樣式設定
        self.title_bg_color = RGBColor(230, 240, 255)  # 主題頁背景（淺藍）
        self.content_bg_color = RGBColor(245, 245, 245)  # 內文頁背景（淺灰）
        self.font_name = "微軟正黑體"  # 可改為 "Arial" 或其他字型
    
    def create_title_slide(self, title_text):
        """建立主題投影片（## 標記）"""
        # 使用空白版面配置
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 設定背景顏色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.title_bg_color
        
        # 新增標題文字框（置中）
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title_text
        
        # 設定標題樣式
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.name = self.font_name
        
        return slide
    
    def create_content_slide(self, title_text):
        """建立內文投影片（# 標記）"""
        # 使用空白版面配置
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 設定背景顏色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.content_bg_color
        
        # 新增標題文字框
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        text_frame = title_box.text_frame
        text_frame.text = title_text
        
        # 設定標題樣式
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.name = self.font_name
        
        # 新增內容文字框
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5)
        )
        
        # 儲存內容文字框供後續使用
        slide.content_textbox = content_box
        
        return slide
    
    def add_content_to_slide(self, slide, content_text):
        """新增內容到投影片"""
        if not hasattr(slide, 'content_textbox'):
            return
        
        text_frame = slide.content_textbox.text_frame
        
        # 新增段落
        if text_frame.text == "":
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = content_text
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = self.font_name
    
    def convert_from_text(self, text_content):
        """從文字內容轉換為 PPT"""
        lines = text_content.split('\n')
        current_slide = None
        
        for line in lines:
            line = line.strip()
            
            if line.startswith('##'):
                # 主題頁面
                title = line[2:].strip()
                current_slide = self.create_title_slide(title)
                
            elif line.startswith('#'):
                # 內文頁面
                title = line[1:].strip()
                current_slide = self.create_content_slide(title)
                
            elif line and current_slide:
                # 新增內容到目前投影片
                self.add_content_to_slide(current_slide, line)
    
    def convert_from_file(self, input_file):
        """從文字檔案轉換"""
        try:
            with open(input_file, 'r', encoding='utf-8') as f:
                text_content = f.read()
            self.convert_from_text(text_content)
        except Exception as e:
            print(f"讀取檔案時發生錯誤: {e}")
            sys.exit(1)
    
    def save(self, output_file):
        """儲存 PPT 檔案"""
        try:
            self.prs.save(output_file)
            print(f"✅ 成功建立 PowerPoint 檔案：{output_file}")
            print(f"📊 總共建立 {len(self.prs.slides)} 張投影片")
        except Exception as e:
            print(f"儲存檔案時發生錯誤: {e}")
            sys.exit(1)


def main():
    """主程式"""
    if len(sys.argv) < 2:
        print("使用方式：")
        print("  python text_to_ppt.py <輸入檔案.txt> [輸出檔案.pptx]")
        print()
        print("範例：")
        print("  python text_to_ppt.py 範例輸入文字.txt")
        print("  python text_to_ppt.py input.txt output.pptx")
        sys.exit(1)
    
    input_file = sys.argv[1]
    
    # 判斷輸出檔案名稱
    if len(sys.argv) >= 3:
        output_file = sys.argv[2]
    else:
        # 自動產生輸出檔名
        base_name = os.path.splitext(input_file)[0]
        output_file = f"{base_name}.pptx"
    
    # 檢查輸入檔案是否存在
    if not os.path.exists(input_file):
        print(f"❌ 錯誤：找不到輸入檔案 '{input_file}'")
        sys.exit(1)
    
    # 執行轉換
    print(f"📝 讀取文字檔案：{input_file}")
    converter = TextToPPTConverter()
    converter.convert_from_file(input_file)
    converter.save(output_file)


if __name__ == "__main__":
    main()
